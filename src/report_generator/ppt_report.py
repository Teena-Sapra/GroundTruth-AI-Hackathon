import os
from typing import Dict, Any, List

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ..utils.logging_utils import get_logger
from ..charts import generate_charts

logger = get_logger(__name__)


def _build_key_highlights(metrics: Dict[str, Any]) -> List[str]:
    overall = metrics["overall"]
    campaign_summary: pd.DataFrame = metrics["campaign_summary"]

    highlights: List[str] = []
    highlights.append(
        f"{int(overall['total_impressions']):,} impressions, "
        f"{int(overall['total_clicks']):,} clicks and "
        f"{int(overall['total_conversions']):,} conversions this week."
    )
    highlights.append(
        f"CTR held at {overall['overall_ctr']:.2%} with an average CPC of "
        f"{overall['overall_cpc']:.2f} and CPA of {overall['overall_cpa']:.2f}."
    )

    if not campaign_summary.empty:
        top_conv = campaign_summary.sort_values("conversions", ascending=False).iloc[0]
        highlights.append(
            f"{top_conv['campaign_id']} was the strongest driver of conversions "
            f"with {int(top_conv['conversions']):,} conversions."
        )

        worst_cpa = campaign_summary.sort_values("cpa", ascending=False).iloc[0]
        highlights.append(
            f"{worst_cpa['campaign_id']} shows the highest CPA at {worst_cpa['cpa']:.2f}, "
            "indicating a focus area for optimization."
        )

    return highlights


def _build_recommendations(metrics: Dict[str, Any]) -> List[str]:
    campaign_summary: pd.DataFrame = metrics["campaign_summary"]

    if campaign_summary.empty:
        return [
            "Increase data coverage and tracking to enable deeper optimization insights.",
            "Introduce multiple creative variants per campaign and measure performance.",
        ]

    recs: List[str] = []

    top_cvr = campaign_summary.sort_values("cvr", ascending=False).head(2)
    low_cpa = campaign_summary.sort_values("cpa", ascending=True).head(2)
    worst_cpa = campaign_summary.sort_values("cpa", ascending=False).head(2)

    top_cvr_campaigns = ", ".join(top_cvr["campaign_id"].astype(str).tolist())
    low_cpa_campaigns = ", ".join(low_cpa["campaign_id"].astype(str).tolist())
    worst_cpa_campaigns = ", ".join(worst_cpa["campaign_id"].astype(str).tolist())

    recs.append(
        f"Prioritize budget allocation towards high-CVR campaigns ({top_cvr_campaigns}) "
        f"and low-CPA campaigns ({low_cpa_campaigns})."
    )
    recs.append(
        f"Review targeting, creatives and landing pages for higher-CPA campaigns "
        f"({worst_cpa_campaigns}) to improve efficiency."
    )
    recs.append(
        "Plan structured A/B tests for creative and audience combinations and consolidate learnings weekly."
    )

    return recs


def _build_top_campaign_rows(campaign_summary: pd.DataFrame) -> List[List[str]]:
    top_campaigns = campaign_summary.sort_values("spend", ascending=False).head(5)
    rows = [["Campaign", "Impr.", "Clicks", "Conv.", "Spend", "CTR", "CPC", "CPA"]]
    for _, row in top_campaigns.iterrows():
        rows.append(
            [
                str(row["campaign_id"]),
                f"{int(row['impressions']):,}",
                f"{int(row['clicks']):,}",
                f"{int(row['conversions']):,}",
                f"{row['spend']:.2f}",
                f"{row['ctr']:.2%}",
                f"{row['cpc']:.2f}",
                f"{row['cpa']:.2f}",
            ]
        )
    return rows


def _detect_anomalies(merged_df: pd.DataFrame, drop_threshold: float = 0.3) -> List[str]:
    """
    Same anomaly logic as PDF: detect big drops in impressions by location vs previous day.
    Returns a list of human-readable strings.
    """
    if "impressions" not in merged_df.columns or "location" not in merged_df.columns or "date" not in merged_df.columns:
        return []

    daily = (
        merged_df.groupby(["date", "location"], as_index=False)["impressions"]
        .sum()
        .sort_values(["location", "date"])
    )

    anomalies: List[str] = []
    has_rain = "rainfall_mm" in merged_df.columns

    for location, grp in daily.groupby("location"):
        grp = grp.sort_values("date")
        for i in range(1, len(grp)):
            prev_row = grp.iloc[i - 1]
            curr_row = grp.iloc[i]
            prev_impr = prev_row["impressions"]
            curr_impr = curr_row["impressions"]
            if prev_impr <= 0:
                continue
            change = (curr_impr - prev_impr) / prev_impr
            if change <= -drop_threshold:
                drop_pct = -change * 100.0
                date_str = str(curr_row["date"])
                msg = (
                    f"Traffic dropped {drop_pct:.1f}% in {location} on {date_str} "
                    f"(impressions {int(curr_impr):,} vs {int(prev_impr):,} previous day)"
                )
                if has_rain:
                    rain_vals = merged_df[
                        (merged_df["date"] == curr_row["date"])
                        & (merged_df["location"] == location)
                    ]["rainfall_mm"]
                    if not rain_vals.empty:
                        rain = float(rain_vals.mean())
                        msg += f". Reported rainfall: {rain:.1f}mm."
                anomalies.append(msg + ".")
    return anomalies


def add_wrapped_text_box(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    font_size: int = 18,
):
    """
    Adds a textbox with word-wrap and safe margins
    so long executive summaries stay inside the slide.
    """
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(30, 30, 30)

    # Simple heuristic for long text
    if len(text) > 900:
        p.font.size = Pt(16)
    if len(text) > 1400:
        p.font.size = Pt(14)

    return textbox


def generate_ppt_report(
    merged_df: pd.DataFrame,
    metrics: Dict[str, Any],
    insights: Dict[str, str],
    config,
    output_path: str,
) -> None:
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = Presentation()

    client_name = config["report"]["client_name"]
    week_start = config["report"]["week_start"]
    week_end = config["report"]["week_end"]

    # Generate charts
    base_output_dir = os.path.dirname(output_path)
    chart_paths = generate_charts(merged_df, metrics, base_output_dir)

    # Layouts
    title_layout = prs.slide_layouts[0]
    title_and_content_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]

    # --- Title slide ---
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = f"{client_name} - Weekly Performance Report"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Reporting Period: {week_start} to {week_end}\nGenerated by Automated Insight Engine"

    # --- Executive Summary slide ---
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = "Executive Summary"

    add_wrapped_text_box(
        slide,
        insights["executive_summary"],
        left=Inches(0.5),
        top=Inches(1.4),
        width=Inches(9.0),
        height=Inches(4.5),
        font_size=18,
    )

    # --- Key Highlights slide ---
    highlights = _build_key_highlights(metrics)
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Key Highlights"

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    if highlights:
        tf.text = highlights[0]
        tf.paragraphs[0].font.size = Pt(22)
        for h in highlights[1:]:
            p = tf.add_paragraph()
            p.text = h
            p.level = 0
            p.font.size = Pt(20)

    # --- KPI Dashboard slide ---
    overall = metrics["overall"]
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "KPI Dashboard"

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    kpi_lines = [
        f"Impressions: {int(overall['total_impressions']):,}",
        f"Clicks: {int(overall['total_clicks']):,}",
        f"Conversions: {int(overall['total_conversions']):,}",
        f"Spend: {overall['total_spend']:.2f}",
        f"CTR: {overall['overall_ctr']:.2%}",
        f"CPC: {overall['overall_cpc']:.2f}",
        f"CVR: {overall['overall_cvr']:.2%}",
        f"CPA: {overall['overall_cpa']:.2f}",
    ]

    tf.text = kpi_lines[0]
    tf.paragraphs[0].font.size = Pt(22)
    for line in kpi_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)

    # --- Top Campaigns Table slide ---
    campaign_summary: pd.DataFrame = metrics["campaign_summary"]
    top_rows = _build_top_campaign_rows(campaign_summary)

    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Top Campaigns (by Spend)"

    rows_count = len(top_rows)
    cols_count = len(top_rows[0])
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(0.8)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    table.columns[0].width = Inches(1.8)
    for col_idx in range(1, cols_count):
        table.columns[col_idx].width = Inches(1.0)

    for r_idx, row_vals in enumerate(top_rows):
        for c_idx, val in enumerate(row_vals):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                if r_idx == 0:
                    paragraph.font.bold = True

    # --- Chart slide: Top Campaigns ---
    if "campaign_bar" in chart_paths:
        slide = prs.slides.add_slide(title_only_layout)
        slide.shapes.title.text = "Top Campaigns – Impressions vs Clicks"

        pic_left = Inches(0.5)
        pic_top = Inches(1.5)
        pic_width = Inches(9.0)
        slide.shapes.add_picture(chart_paths["campaign_bar"], pic_left, pic_top, width=pic_width)

    # --- Chart slide: Daily Trend ---
    if "daily_trend" in chart_paths:
        slide = prs.slides.add_slide(title_only_layout)
        slide.shapes.title.text = "Daily Trend – Impressions & Clicks"

        pic_left = Inches(0.5)
        pic_top = Inches(1.5)
        pic_width = Inches(9.0)
        slide.shapes.add_picture(chart_paths["daily_trend"], pic_left, pic_top, width=pic_width)

    # --- Anomaly Analysis slide ---
    anomalies = _detect_anomalies(merged_df)
    if anomalies:
        slide = prs.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = "Anomaly Analysis"

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        tf.text = anomalies[0]
        tf.paragraphs[0].font.size = Pt(20)
        for a in anomalies[1:]:
            p = tf.add_paragraph()
            p.text = a
            p.level = 0
            p.font.size = Pt(18)

    # --- Optimization Plan slide ---
    recs = _build_recommendations(metrics)
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Optimization Plan"

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    if recs:
        tf.text = recs[0]
        tf.paragraphs[0].font.size = Pt(22)
        for r in recs[1:]:
            p = tf.add_paragraph()
            p.text = r
            p.level = 0
            p.font.size = Pt(20)

    # --- Glossary slide ---
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Glossary"

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    glossary_items = [
        "CTR (Click-Through Rate): Percentage of impressions that resulted in clicks.",
        "CVR (Conversion Rate): Percentage of clicks that resulted in conversions.",
        "CPA (Cost Per Acquisition): Average cost required to generate a single conversion."
    ]

    tf.text = glossary_items[0]
    tf.paragraphs[0].font.size = Pt(20)
    for g in glossary_items[1:]:
        p = tf.add_paragraph()
        p.text = g
        p.level = 0
        p.font.size = Pt(18)

    # --- Closing slide ---
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Thank You"

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    tf.text = f"Report generated by the Automated Insight Engine for {client_name}."
    tf.paragraphs[0].font.size = Pt(22)
    p = tf.add_paragraph()
    p.text = "For questions or further deep dives, please reach out to your account team."
    p.level = 0
    p.font.size = Pt(20)

    prs.save(output_path)
    logger.info("PPTX report generated at %s", output_path)
